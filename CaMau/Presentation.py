from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.util import Inches


# Tạo một Presentation mới
prs = Presentation()
# Thêm một slide đầu tiên (Title slide hoặc slide bất kỳ)
slide_layout = prs.slide_layouts[0]  # Title slide
slide = prs.slides.add_slide(slide_layout)

# Thêm hình ảnh vào slide đầu tiên
slide.shapes.add_picture("anh.jpg", Inches(1), Inches(1), Inches(8), Inches(4))

# Hàm hỗ trợ để thêm slide với title và bullet points
def add_slide_with_bullets(prs, title_text, bullets, subtitle=None, quote=None):
    slide_layout = prs.slide_layouts[1]  # Layout với title và content
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = title_text
    title.text_frame.paragraphs[0].font.size = Pt(32)
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 128, 0)  # Màu xanh

    if subtitle:
        tf = title.text_frame
        p = tf.add_paragraph()
        p.text = subtitle
        p.font.size = Pt(18)
        p.font.italic = True

    body_shape = slide.placeholders[1]
    tf = body_shape.text_frame
    tf.clear()  # Xóa nội dung mặc định
    for bullet in bullets:
        p = tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(14)
        p.level = 1 if bullet.startswith("  -") else 0  # Hỗ trợ bullet con
        p.font.color.rgb = RGBColor(0, 0, 0)  # Màu chữ đen

    if quote:
        p = tf.add_paragraph()
        p.text = quote
        p.font.italic = True
        p.font.size = Pt(12)
        p.alignment = PP_ALIGN.CENTER


# Hàm hỗ trợ để thêm slide với table (cho slide 6)
def add_slide_with_table(prs, title_text, table_data):
    slide_layout = prs.slide_layouts[5]  # Layout với title only
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = title_text
    title.text_frame.paragraphs[0].font.size = Pt(32)
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 128, 0)

    # Thêm table
    rows, cols = len(table_data), len(table_data[0])
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(9)
    height = Inches(4.5)
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table

    # Đặt header
    for col in range(cols):
        cell = table.cell(0, col)
        cell.text = table_data[0][col]
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0, 128, 0)  # Màu xanh
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(12)

    # Đổ dữ liệu
    for row in range(1, rows):
        for col in range(cols):
            cell = table.cell(row, col)
            cell.text = table_data[row][col]
            cell.text_frame.paragraphs[0].font.size = Pt(10)


# Slide 1: Tiêu đề
slide_layout = prs.slide_layouts[0]  # Title slide
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Giải pháp Bản đồ số Nông nghiệp cho Sở NN&PTNT Cà Mau"
title.text_frame.paragraphs[0].font.size = Pt(36)
title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 128, 0)
subtitle = slide.placeholders[1]
subtitle.text = "Tích hợp viễn thám, IoT và GIS để tối ưu hóa quản lý tài nguyên, nâng cao hiệu quả sản xuất, bảo vệ môi trường và hỗ trợ chuyển đổi số"
subtitle.text_frame.paragraphs[0].font.size = Pt(20)
subtitle.text_frame.paragraphs[0].font.italic = True

# Slide 2: Căn cứ Pháp lý và Chính sách
bullets = [
    "- Nghị quyết số 19-NQ/TW: Phát triển nông nghiệp thông minh, chuyển đổi số ngành nông nghiệp.",
    "- Chỉ thị số 15/CT-TTg: Khuyến khích sử dụng công nghệ cao, bao gồm bản đồ số, viễn thám và IoT.",
    "- Kế hoạch số 429/KH-BNN-CN: Xây dựng kế hoạch chuyển đổi số ngành nông nghiệp, nhấn mạnh tích hợp GIS và cảm biến.",
    "- Tiêu chuẩn kỹ thuật quốc tế: ISO 19152:2012 (LADM cho quản lý đất đai), ISO 19115:2003 (Metadata cho thông tin địa lý), và các tiêu chuẩn IoT như IEEE 802.15.4 cho cảm biến không dây."
]
quote = '"Việc triển khai bản đồ số nông nghiệp, kết hợp viễn thám và IoT, phù hợp với chiến lược phát triển nông nghiệp thông minh và bền vững của Chính phủ, giúp nâng cao hiệu quả quản lý tại Cà Mau."'
add_slide_with_bullets(prs, "Căn cứ Pháp lý và Chính sách hỗ trợ", bullets, quote=quote)

# Slide 3: Hiện trạng Nông nghiệp tại Cà Mau
bullets = [
    "- Ngành nông nghiệp đóng vai trò quan trọng trong GRDP của tỉnh (31,9%), với trọng tâm là nuôi trồng thủy sản và trồng trọt.",
    "- Nuôi trồng thủy sản: Diện tích nuôi tôm: 427.000 ha, chiếm 45% diện tích nuôi tôm của Đồng bằng Sông Cửu Long; mô hình nuôi tôm sạch RAS đang phát triển mạnh, nhưng thiếu giám sát chất lượng nước thời gian thực (như Chlorophyll-a, độ đục, nhiệt độ).",
    "- Trồng trọt và chăn nuôi: Mô hình lúa – tôm: Giúp tăng năng suất và bảo vệ môi trường; các vùng trồng lúa, trang trại và công trình thủy lợi cần theo dõi biến động đất đai qua viễn thám.",
    "- Chuyển đổi số: Các phần mềm hỗ trợ quản lý nông nghiệp và thủy sản đã được triển khai, nhưng chưa tích hợp đầy đủ IoT (cảm biến, camera) và viễn thám để giám sát 24/7, cảnh báo sớm dịch bệnh hoặc biến đổi môi trường."
]
quote = '"Cà Mau đang trong quá trình chuyển đổi số ngành nông nghiệp, đặc biệt trong quản lý tài nguyên và sản xuất thủy sản, nhưng cần hệ thống tích hợp để theo dõi liên tục và dự báo chính xác hơn."'
add_slide_with_bullets(prs, "Hiện trạng Nông nghiệp tại Cà Mau", bullets, quote=quote)

# Slide 4: Các Vấn đề và Thách thức Hiện tại
bullets = [
    "- Quản lý tài nguyên chưa tối ưu: Giám sát diện tích đất, chất lượng đất và nước chưa liên tục; thiếu dữ liệu thời gian thực từ cảm biến (độ ẩm, pH, oxy hòa tan) dẫn đến lãng phí tài nguyên.",
    "- Ô nhiễm môi trường: Nuôi trồng thủy sản gặp ô nhiễm nước (tảo nở hoa, phú dưỡng hóa); trồng trọt thiếu cảnh báo biến đổi môi trường định kỳ qua viễn thám.",
    "- Thiếu kết nối và minh bạch: Quản lý sản phẩm nông sản/thủy sản (truy xuất nguồn gốc) chưa tích hợp blockchain; thiếu giám sát sức khỏe vật nuôi/cây trồng qua camera đa quang phổ và AI.",
    "- Khả năng dự báo và cảnh báo: Chưa có hệ thống cảnh báo sớm về thời tiết cực đoan, dịch bệnh (dịch bệnh tôm, cá, lúa) hoặc biến động rừng ngập mặn, dẫn đến thiệt hại lớn."
]
quote = '"Mặc dù đã có tiến triển trong chuyển đổi số, nhưng vẫn thiếu hệ thống quản lý toàn diện với viễn thám, IoT và AI để tối ưu hóa sản xuất, giảm rủi ro và bảo vệ môi trường tại Cà Mau."'
add_slide_with_bullets(prs, "Các Vấn đề và Thách thức trong Quản lý Nông nghiệp", bullets, quote=quote)

# Slide 5: Giải pháp Bản đồ số Nông nghiệp
bullets = [
    "- Quản lý tài nguyên nông nghiệp và thủy sản:",
    "  - Phân hệ 1: Quản lý vùng nuôi thủy sản (tôm sú, tôm càng xanh, cá tra, cá lóc, cá rô phi) qua bản đồ GIS, viễn thám để lập bản đồ vùng nuôi, giám sát chất lượng nước (Chlorophyll-a, độ đục, nhiệt độ) và dự báo môi trường.",
    "  - Phân hệ 2: Quản lý nông nghiệp (vùng trồng lúa, trang trại, công trình thủy lợi, lò giết mổ, PCCC rừng) với dữ liệu GIS, cảm biến và AI để theo dõi biến động đất đai, sức khỏe cây trồng/vật nuôi.",
    "- Ứng dụng viễn thám và GIS: Giám sát thay đổi diện tích đất, chất lượng nước/môi trường thời gian thực; cảnh báo sớm thời tiết cực đoan và biến đổi môi trường.",
    "- Tích hợp IoT (cảm biến, camera): Giám sát 24/7 chất lượng nước (pH, oxy, độ mặn), sức khỏe thủy sản/cây trồng (phát hiện sâu bệnh qua camera đa quang phổ, UAV); tự động hóa tưới tiêu, cho ăn.",
    "- AI và Big Data: Phân tích dữ liệu từ cảm biến/viễn thám để dự báo dịch bệnh, sản lượng; quản lý nhật ký truy cập, thống kê sản lượng theo loài/phường xã.",
    "- Ứng dụng mobile: Cho chủ hộ nuôi trồng cập nhật dịch bệnh, sản lượng; xem bản đồ cảnh báo."
]
quote = '"Giải pháp bản đồ số tích hợp viễn thám, IoT và AI sẽ giúp Cà Mau chuyển đổi ngành nông nghiệp thành hệ thống thông minh, giảm thiểu rủi ro và tối ưu hóa tài nguyên bền vững."'
add_slide_with_bullets(prs, "Giải pháp Bản đồ số Nông nghiệp cho Cà Mau", bullets, quote=quote)

# Slide 6: Các Lớp Dữ Liệu (CSDL) cho Giải pháp
table_data = [
    ["Lớp dữ liệu", "Mô tả"],
    ["Vùng nuôi thủy sản",
     "Diện tích nuôi tôm sú, tôm càng xanh, cá tra, cá lóc, cá rô phi; theo dõi biến động qua viễn thám."],
    ["Công trình thủy lợi", "Kênh mương, đê điều, hệ thống cấp thoát nước; lọc theo năm xây dựng."],
    ["Trang trại và mô hình lúa – tôm", "Diện tích, năng suất, chất lượng đất; lọc theo thời gian xây dựng."],
    ["Lò giết mổ động vật", "Quản lý chất lượng sản phẩm; lọc theo năm mở."],
    ["Môi trường nuôi trồng",
     "Chất lượng nước (độ mặn, pH, nhiệt độ, oxy hòa tan từ cảm biến); rừng ngập mặn để giám sát biến động."],
    ["Sự biến động môi trường", "Biến đổi khí hậu, thời tiết cực đoan; cảnh báo dịch bệnh (tôm, cá, lúa)."],
    ["Truy xuất nguồn gốc", "Blockchain cho sản phẩm nông sản/thủy sản."],
    ["Sản phẩm và thống kê", "Sản lượng, xuất khẩu, chuỗi cung ứng; dashboard so sánh sản lượng theo phường/xã."],
    ["Hình ảnh UAV", "Quản lý đợt quay chụp để giám sát môi trường và vùng nuôi."],
    ["Quản trị hệ thống", "Người dùng, vai trò, nhật ký truy cập/thay đổi dữ liệu."]
]
add_slide_with_table(prs, "Các lớp dữ liệu quan trọng cho Cà Mau", table_data)

# Slide 7: Công nghệ và Phương Pháp Tiến Hành
bullets = [
    "- Viễn thám: Ảnh vệ tinh độ phân giải cao để phân tích biến động đất đai, chất lượng nước (Chlorophyll-a, độ đục); so sánh định kỳ để dự báo môi trường.",
    "- GIS: Quản lý dữ liệu không gian (đa giác, điểm); công cụ tìm kiếm thuộc tính/không gian, đo đạc, chỉ đường; import dữ liệu (Excel, GeoJSON, Shapefile).",
    "- IoT (Cảm biến và Camera): Giám sát thời gian thực (độ ẩm đất, nhiệt độ, hành vi vật nuôi qua camera dưới nước/multispectral); tích hợp AI để cảnh báo sớm (dịch bệnh, tảo nở hoa), tự động hóa (tưới tiêu, cho ăn).",
    "- AI và Big Data: Phân tích dữ liệu lớn từ đám mây để dự báo tăng trưởng, dịch bệnh; dashboard thống kê sản lượng, biểu đồ biến động theo năm.",
    "- Ứng dụng mobile: Đăng nhập, cập nhật dịch bệnh/sản lượng, xem bản đồ cảnh báo.",
    "- Triển khai: Sử dụng môi trường Python cho AI (với thư viện như numpy, scipy cho phân tích); tích hợp UAV cho hình ảnh."
]
add_slide_with_bullets(prs, "Công nghệ và Phương pháp Triển khai Giải pháp", bullets)

# Slide 8: Lợi ích và Kết quả Dự Kiến
bullets = [
    "- Quản lý tài nguyên hiệu quả: Giám sát liên tục diện tích đất/nước qua viễn thám/IoT, tối ưu hóa đầu vào (tiết kiệm nước 20-30% qua tưới tiêu chính xác).",
    "- Tăng năng suất và chất lượng: Phát hiện sớm sâu bệnh/dịch bệnh qua camera/AI, tăng sản lượng tôm/lúa 15-25%; cải thiện chất lượng sản phẩm cho xuất khẩu.",
    "- Bảo vệ môi trường: Giám sát chất lượng nước/đất, bảo vệ rừng ngập mặn; giảm ô nhiễm từ nuôi trồng (giảm lãng phí thức ăn qua cho ăn tự động).",
    "- Phát triển bền vững: Thống kê báo cáo chi tiết (Excel xuất), minh bạch truy xuất nguồn gốc; hỗ trợ nông dân qua ứng dụng mobile, giảm thiệt hại từ biến đổi môi trường lên đến 40%.",
    "- Kết quả kinh tế: Nâng cao GRDP nông nghiệp, hỗ trợ mô hình RAS và lúa-tôm bền vững."
]
add_slide_with_bullets(prs, "Lợi ích và Kết quả Dự Kiến", bullets)

# Slide 9: Đề xuất Hành động cho Sở NN&PTNT Cà Mau
bullets = [
    "- Khởi động dự án thí điểm: Tại vùng nuôi tôm sú/cá tra và mô hình lúa-tôm để kiểm tra viễn thám/IoT, điều chỉnh dựa trên dữ liệu thực tế.",
    "- Hợp tác đối tác: Với công ty GIS (ESRI/ArcGIS), viễn thám (Sentinel/Landsat), IoT (cảm biến từ Arduino/Raspberry Pi) và AI (Google Cloud/TensorFlow).",
    "- Đào tạo: Cán bộ về quản lý hệ thống (quyền vai trò, import dữ liệu); nông dân về ứng dụng mobile và cảm biến để cập nhật dịch bệnh/sản lượng.",
    "- Lộ trình triển khai: Giai đoạn 1 (6 tháng): Xây dựng lớp dữ liệu cơ bản và dashboard. Giai đoạn 2 (1 năm): Tích hợp IoT/AI. Giai đoạn 3 (2-3 năm): Mở rộng toàn tỉnh, với ngân sách ước tính và đo lường KPI (tăng năng suất, giảm rủi ro)."
]
add_slide_with_bullets(prs, "Đề xuất Hành động Triển khai Giải pháp", bullets)

# Slide 10: Kết luận
bullets = [
    "- Giải pháp bản đồ số nông nghiệp, tích hợp viễn thám, IoT và AI, sẽ giúp Sở NN&PTNT Cà Mau quản lý tài nguyên hiệu quả, nâng cao năng suất/chất lượng sản phẩm, bảo vệ môi trường và hướng tới phát triển bền vững trong nông nghiệp/thủy sản.",
    "- Sự phối hợp giữa cơ quan nhà nước, đối tác công nghệ và cộng đồng nông dân là chìa khóa thành công, mang lại lợi ích lâu dài cho tỉnh Cà Mau."
]
quote = '"Giải pháp bản đồ số nông nghiệp không chỉ giúp Cà Mau đạt mục tiêu phát triển bền vững mà còn nâng tầm hiệu quả quản lý và sản xuất, phù hợp với chuyển đổi số quốc gia."'
add_slide_with_bullets(prs, "Kết luận", bullets, quote=quote)

# Lưu file PPTX
prs.save("giai_phap_ban_do_so_nong_nghiep.pptx")
print("File PPTX đã được tạo: giai_phap_ban_do_so_nong_nghiep.pptx")